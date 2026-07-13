#!/usr/bin/env python3
"""Create deterministic, reviewable fixtures for the CTOX Office oracle."""

from pathlib import Path
from datetime import datetime, timezone
from io import BytesIO
import re
import struct
import zlib
from zipfile import ZipFile, ZipInfo, ZIP_DEFLATED
from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches as PptxInches
from lxml import etree


ROOT = Path(__file__).resolve().parents[3]
OUTPUT_DIRECTORY = ROOT / "tests/fixtures/office/document"
OPEN_RENDER_OUTPUT = OUTPUT_DIRECTORY / "open-render-zoom.docx"
EDIT_SAVE_OUTPUT = OUTPUT_DIRECTORY / "edit-save.docx"
UNDO_CLIPBOARD_OUTPUT = OUTPUT_DIRECTORY / "undo-clipboard-keyboard.docx"
CHARACTER_PARAGRAPH_OUTPUT = OUTPUT_DIRECTORY / "character-paragraph-formatting.docx"
STYLES_LISTS_NUMBERING_OUTPUT = OUTPUT_DIRECTORY / "styles-lists-numbering.docx"
TABLES_OUTPUT = OUTPUT_DIRECTORY / "tables.docx"
IMAGES_POSITIONING_OUTPUT = OUTPUT_DIRECTORY / "images-positioning.docx"
SECTIONS_HEADERS_FOOTERS_OUTPUT = OUTPUT_DIRECTORY / "sections-headers-footers.docx"
LINKS_BOOKMARKS_FIELDS_OUTPUT = OUTPUT_DIRECTORY / "links-bookmarks-fields.docx"
COMMENTS_TRACK_CHANGES_OUTPUT = OUTPUT_DIRECTORY / "comments-track-changes.docx"
DRAWINGS_CHARTS_OUTPUT = OUTPUT_DIRECTORY / "drawings-charts.docx"
SPREADSHEET_OUTPUT_DIRECTORY = ROOT / "tests/fixtures/office/spreadsheet"
SPREADSHEET_OPEN_RENDER_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "open-render-sheets.xlsx"
SPREADSHEET_EDIT_SAVE_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "edit-save.xlsx"
SPREADSHEET_UNDO_CLIPBOARD_FILL_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "undo-clipboard-fill.xlsx"
SPREADSHEET_CELL_FORMAT_ROWS_COLUMNS_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "cell-format-rows-columns.xlsx"
SPREADSHEET_FORMULAS_REFERENCES_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "formulas-references.xlsx"
SPREADSHEET_MULTI_SHEET_MERGE_FREEZE_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "multi-sheet-merge-freeze.xlsx"
SPREADSHEET_SORT_FILTER_TABLES_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "sort-filter-tables.xlsx"
SPREADSHEET_VALIDATION_CONDITIONAL_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "validation-conditional-formatting.xlsx"
SPREADSHEET_COMMENTS_NAMES_PROTECTION_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "comments-names-protection.xlsx"
SPREADSHEET_CHARTS_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "charts.xlsx"
SPREADSHEET_PIVOT_PRINT_LAYOUT_OUTPUT = SPREADSHEET_OUTPUT_DIRECTORY / "pivot-print-layout.xlsx"


def write_solid_png(path: Path, width: int, height: int, rgb: tuple[int, int, int]) -> None:
    """Write a tiny deterministic RGB PNG without adding an image dependency."""
    signature = b"\x89PNG\r\n\x1a\n"

    def chunk(kind: bytes, payload: bytes) -> bytes:
        return (
            struct.pack(">I", len(payload))
            + kind
            + payload
            + struct.pack(">I", zlib.crc32(kind + payload) & 0xFFFFFFFF)
        )

    scanline = bytes(rgb) * width
    pixels = b"".join(bytes([0]) + scanline for _ in range(height))
    path.write_bytes(
        signature
        + chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(pixels, 9))
        + chunk(b"IEND", b"")
    )


def make_floating_picture(inline, horizontal_twips: int, vertical_twips: int) -> None:
    """Turn python-docx's inline DrawingML into a square-wrapped page anchor."""
    inline.tag = qn("wp:anchor")
    for name, value in {
        "distT": "0", "distB": "0", "distL": "114300", "distR": "114300",
        "simplePos": "0", "relativeHeight": "251658240", "behindDoc": "0",
        "locked": "0", "layoutInCell": "1", "allowOverlap": "1",
    }.items():
        inline.set(name, value)
    extent = inline.find(qn("wp:extent"))
    simple_position = OxmlElement("wp:simplePos")
    simple_position.set("x", "0")
    simple_position.set("y", "0")
    horizontal = OxmlElement("wp:positionH")
    horizontal.set("relativeFrom", "column")
    horizontal_offset = OxmlElement("wp:posOffset")
    horizontal_offset.text = str(horizontal_twips * 635)
    horizontal.append(horizontal_offset)
    vertical = OxmlElement("wp:positionV")
    vertical.set("relativeFrom", "paragraph")
    vertical_offset = OxmlElement("wp:posOffset")
    vertical_offset.text = str(vertical_twips * 635)
    vertical.append(vertical_offset)
    inline.insert(0, simple_position)
    inline.insert(1, horizontal)
    inline.insert(2, vertical)
    wrap = OxmlElement("wp:wrapSquare")
    wrap.set("wrapText", "bothSides")
    inline.insert(inline.index(extent) + 1, wrap)


def add_external_hyperlink(paragraph, text: str, url: str) -> None:
    relationship_id = paragraph.part.relate_to(url, RT.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)
    run = OxmlElement("w:r")
    properties = OxmlElement("w:rPr")
    style = OxmlElement("w:rStyle")
    style.set(qn("w:val"), "Hyperlink")
    properties.append(style)
    run.append(properties)
    value = OxmlElement("w:t")
    value.text = text
    run.append(value)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def add_bookmark(paragraph, name: str, bookmark_id: int, text: str) -> None:
    start = OxmlElement("w:bookmarkStart")
    start.set(qn("w:id"), str(bookmark_id))
    start.set(qn("w:name"), name)
    paragraph._p.append(start)
    paragraph.add_run(text)
    end = OxmlElement("w:bookmarkEnd")
    end.set(qn("w:id"), str(bookmark_id))
    paragraph._p.append(end)


def add_complex_field(paragraph, instruction: str, cached_text: str) -> None:
    begin_run = OxmlElement("w:r")
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    begin.set(qn("w:dirty"), "true")
    begin_run.append(begin)
    paragraph._p.append(begin_run)
    instruction_run = OxmlElement("w:r")
    instruction_text = OxmlElement("w:instrText")
    instruction_text.set(qn("xml:space"), "preserve")
    instruction_text.text = f" {instruction} "
    instruction_run.append(instruction_text)
    paragraph._p.append(instruction_run)
    separate_run = OxmlElement("w:r")
    separate = OxmlElement("w:fldChar")
    separate.set(qn("w:fldCharType"), "separate")
    separate_run.append(separate)
    paragraph._p.append(separate_run)
    paragraph.add_run(cached_text)
    end_run = OxmlElement("w:r")
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    end_run.append(end)
    paragraph._p.append(end_run)


def add_comment_range(paragraph, text: str, comment_id: int) -> None:
    run = paragraph.add_run(text)
    start = OxmlElement("w:commentRangeStart")
    start.set(qn("w:id"), str(comment_id))
    run._r.addprevious(start)
    end = OxmlElement("w:commentRangeEnd")
    end.set(qn("w:id"), str(comment_id))
    run._r.addnext(end)
    reference_run = OxmlElement("w:r")
    properties = OxmlElement("w:rPr")
    style = OxmlElement("w:rStyle")
    style.set(qn("w:val"), "CommentReference")
    properties.append(style)
    reference_run.append(properties)
    reference = OxmlElement("w:commentReference")
    reference.set(qn("w:id"), str(comment_id))
    reference_run.append(reference)
    end.addnext(reference_run)


def add_existing_revision(paragraph, kind: str, revision_id: int, text: str) -> None:
    revision = OxmlElement(f"w:{kind}")
    revision.set(qn("w:id"), str(revision_id))
    revision.set(qn("w:author"), "CTOX Existing Reviewer")
    revision.set(qn("w:date"), "2026-07-11T00:00:00Z")
    run = OxmlElement("w:r")
    value = OxmlElement("w:t" if kind == "ins" else "w:delText")
    value.text = text
    run.append(value)
    revision.append(run)
    paragraph._p.append(revision)


def set_cell_shading(cell, fill: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    properties.append(shading)


def add_page(document: Document, number: int, accent: str, label: str) -> None:
    title = document.add_paragraph()
    title.style = document.styles["Title"]
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(f"ORACLE PAGE {number}")
    run.font.name = "Arial"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(accent)

    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(label)
    run.font.name = "Arial"
    run.font.size = Pt(16)
    run.font.bold = True

    body = document.add_paragraph()
    body.paragraph_format.space_before = Pt(18)
    body.paragraph_format.space_after = Pt(12)
    run = body.add_run(
        "Deterministische CTOX-Office-Fixture für Seitenrendering, Navigation und Zoom. "
        "Die sichtbare Seitennummer, Farbe und Prüflinie unterscheiden jede Seite eindeutig."
    )
    run.font.name = "Arial"
    run.font.size = Pt(12)

    table = document.add_table(rows=2, cols=2)
    table.autofit = False
    widths = [Inches(2.0), Inches(4.2)]
    values = [("Feature", "document.open-render-zoom"), ("Prüfwert", f"CTOX-{number:02d}-{accent}")]
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            cell.width = widths[column_index]
            cell.text = values[row_index][column_index]
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_after = Pt(3)
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(10)
            if column_index == 0:
                set_cell_shading(cell, "E7EFED")
                cell.paragraphs[0].runs[0].font.bold = True

    marker = document.add_paragraph()
    marker.paragraph_format.space_before = Pt(24)
    marker.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = marker.add_run(f"END OF PAGE {number} — {label}")
    run.font.name = "Arial"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(accent)


def create_open_render_zoom() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    section.header_distance = Inches(0.5)
    section.footer_distance = Inches(0.5)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(12)
    normal.paragraph_format.space_after = Pt(8)
    document.core_properties.title = "CTOX Oracle Open Render Zoom Fixture"
    document.core_properties.subject = "document.open-render-zoom"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 10, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    pages = [(1, "176B5B", "GREEN"), (2, "1F5A92", "BLUE"), (3, "8B3A62", "MAGENTA")]
    for index, (number, accent, label) in enumerate(pages):
        if index:
            document.add_page_break()
        add_page(document, number, accent, label)

    OPEN_RENDER_OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    document.save(OPEN_RENDER_OUTPUT)
    canonicalize_zip(OPEN_RENDER_OUTPUT)
    print(OPEN_RENDER_OUTPUT)


def create_edit_save() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(12)
    document.core_properties.title = "CTOX Oracle Edit Save Fixture"
    document.core_properties.subject = "document.edit-save"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 10, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX EDIT / SAVE ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("176B5B")

    instruction = document.add_paragraph()
    instruction.paragraph_format.space_before = Pt(24)
    instruction.add_run("Replace the complete marker on the next line and save the document.")

    target = document.add_paragraph()
    target.paragraph_format.space_before = Pt(18)
    run = target.add_run("CTOX_EDIT_TARGET_ALPHA")
    run.font.name = "Arial"
    run.font.size = Pt(16)
    run.font.bold = True

    footer = document.add_paragraph()
    footer.paragraph_format.space_before = Pt(36)
    footer.add_run("PRESERVE_THIS_UNRELATED_PARAGRAPH_7B21")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(EDIT_SAVE_OUTPUT)
    canonicalize_zip(EDIT_SAVE_OUTPUT)
    print(EDIT_SAVE_OUTPUT)


def create_undo_clipboard_keyboard() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(12)
    document.core_properties.title = "CTOX Oracle Undo Clipboard Keyboard Fixture"
    document.core_properties.subject = "document.undo-clipboard-keyboard"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 10, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX UNDO / CLIPBOARD ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("1F5A92")

    instruction = document.add_paragraph()
    instruction.paragraph_format.space_before = Pt(20)
    instruction.add_run("Exercise keyboard input, undo, redo, copy, cut, and paste.")

    source = document.add_paragraph()
    source.paragraph_format.space_before = Pt(18)
    run = source.add_run("UNDO_CLIPBOARD_BASE")
    run.font.name = "Arial"
    run.font.size = Pt(16)
    run.font.bold = True

    destination = document.add_paragraph()
    destination.paragraph_format.space_before = Pt(18)
    run = destination.add_run("PASTE_DESTINATION")
    run.font.name = "Arial"
    run.font.size = Pt(16)

    footer = document.add_paragraph()
    footer.paragraph_format.space_before = Pt(30)
    footer.add_run("PRESERVE_UNDO_CLIPBOARD_UNRELATED_4C19")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(UNDO_CLIPBOARD_OUTPUT)
    canonicalize_zip(UNDO_CLIPBOARD_OUTPUT)
    print(UNDO_CLIPBOARD_OUTPUT)


def create_character_paragraph_formatting() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(5)
    document.core_properties.title = "CTOX Oracle Character Paragraph Formatting Fixture"
    document.core_properties.subject = "document.character-paragraph-formatting"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 10, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX CHARACTER / PARAGRAPH FORMAT ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("7A3E00")

    instruction = document.add_paragraph()
    instruction.add_run(
        "Format each marker independently; unchanged labels and the preservation marker are controls."
    )

    character_heading = document.add_paragraph()
    character_heading.paragraph_format.space_before = Pt(8)
    heading_run = character_heading.add_run("CHARACTER TARGETS")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("7A3E00")

    for marker in (
        "CHAR_BOLD_TARGET",
        "CHAR_ITALIC_TARGET",
        "CHAR_UNDERLINE_TARGET",
        "CHAR_SIZE_TARGET",
        "CHAR_COLOR_TARGET",
    ):
        paragraph = document.add_paragraph()
        run = paragraph.add_run(marker)
        run.font.name = "Arial"
        run.font.size = Pt(11)

    paragraph_heading = document.add_paragraph()
    paragraph_heading.paragraph_format.space_before = Pt(10)
    heading_run = paragraph_heading.add_run("PARAGRAPH TARGETS")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("7A3E00")

    for marker in (
        "PARA_ALIGN_CENTER_TARGET",
        "PARA_INDENT_TARGET",
        "PARA_LINE_SPACING_TARGET",
    ):
        paragraph = document.add_paragraph()
        run = paragraph.add_run(marker)
        run.font.name = "Arial"
        run.font.size = Pt(11)

    control = document.add_paragraph()
    control.paragraph_format.space_before = Pt(12)
    control.add_run("PRESERVE_FORMATTING_UNRELATED_91D7")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(CHARACTER_PARAGRAPH_OUTPUT)
    canonicalize_zip(CHARACTER_PARAGRAPH_OUTPUT)
    print(CHARACTER_PARAGRAPH_OUTPUT)


def create_styles_lists_numbering() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(5)
    document.core_properties.title = "CTOX Oracle Styles Lists Numbering Fixture"
    document.core_properties.subject = "document.styles-lists-numbering"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 11, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX STYLES / LISTS / NUMBERING ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("385723")

    instruction = document.add_paragraph()
    instruction.add_run(
        "Apply paragraph styles, change a bullet level, and verify numbered-list continuation."
    )

    heading = document.add_paragraph()
    heading.paragraph_format.space_before = Pt(8)
    heading_run = heading.add_run("STYLE TARGETS")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("385723")

    document.add_paragraph("STYLE_HEADING1_TARGET")
    document.add_paragraph("STYLE_QUOTE_TARGET")

    list_heading = document.add_paragraph()
    list_heading.paragraph_format.space_before = Pt(8)
    heading_run = list_heading.add_run("BULLET TARGETS")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("385723")

    document.add_paragraph("BULLET_BASE_ALPHA", style="List Bullet")
    document.add_paragraph("BULLET_NEST_TARGET", style="List Bullet")

    number_heading = document.add_paragraph()
    number_heading.paragraph_format.space_before = Pt(8)
    heading_run = number_heading.add_run("NUMBERING TARGETS")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("385723")

    document.add_paragraph("NUMBER_BASE_ONE", style="List Number")
    document.add_paragraph("NUMBER_BASE_TWO", style="List Number")
    document.add_paragraph("NUMBER_BASE_THREE", style="List Number")

    continuation_heading = document.add_paragraph()
    continuation_heading.paragraph_format.space_before = Pt(8)
    continuation_heading.add_run("CONTINUATION CONTROL")
    document.add_paragraph("NUMBER_CONTINUE_TARGET")

    control = document.add_paragraph()
    control.paragraph_format.space_before = Pt(10)
    control.add_run("PRESERVE_STYLES_LISTS_UNRELATED_E52A")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(STYLES_LISTS_NUMBERING_OUTPUT)
    canonicalize_zip(STYLES_LISTS_NUMBERING_OUTPUT)
    print(STYLES_LISTS_NUMBERING_OUTPUT)


def create_tables() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10)
    normal.paragraph_format.space_after = Pt(4)
    document.core_properties.title = "CTOX Oracle Tables Fixture"
    document.core_properties.subject = "document.tables"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 11, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX TABLES ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("1F5A92")

    instruction = document.add_paragraph(
        "Edit cells, add a row and column, merge and split cells, and preserve the nested table."
    )
    instruction.paragraph_format.space_after = Pt(8)

    main_heading = document.add_paragraph()
    heading_run = main_heading.add_run("MAIN TABLE TARGETS")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("1F5A92")

    table = document.add_table(rows=3, cols=3)
    table.style = "Table Grid"
    table.autofit = False
    widths = [Inches(1.8), Inches(2.2), Inches(2.2)]
    values = [
        ["TABLE_HEADER_A", "TABLE_HEADER_B", "TABLE_HEADER_C"],
        ["TABLE_EDIT_TARGET", "TABLE_ROW_ANCHOR", "TABLE_COLUMN_ANCHOR"],
        ["TABLE_KEEP_A", "TABLE_KEEP_B", "TABLE_KEEP_C"],
    ]
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            cell.width = widths[column_index]
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            cell.text = values[row_index][column_index]
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_after = Pt(0)
                for cell_run in paragraph.runs:
                    cell_run.font.name = "Arial"
                    cell_run.font.size = Pt(9)
                    cell_run.font.bold = row_index == 0
            if row_index == 0:
                set_cell_shading(cell, "D9EAF7")

    merge_heading = document.add_paragraph()
    merge_heading.paragraph_format.space_before = Pt(8)
    heading_run = merge_heading.add_run("MERGE / SPLIT TARGET")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("1F5A92")

    merge_table = document.add_table(rows=2, cols=2)
    merge_table.style = "Table Grid"
    merge_table.autofit = False
    merge_values = [
        ["TABLE_MERGE_LEFT", "TABLE_MERGE_RIGHT"],
        ["TABLE_SPLIT_TARGET", "TABLE_SPLIT_CONTROL"],
    ]
    for row_index, row in enumerate(merge_table.rows):
        for column_index, cell in enumerate(row.cells):
            cell.width = Inches(3.1)
            cell.text = merge_values[row_index][column_index]
            cell.paragraphs[0].paragraph_format.space_after = Pt(0)
            cell.paragraphs[0].runs[0].font.name = "Arial"
            cell.paragraphs[0].runs[0].font.size = Pt(9)

    nested_heading = document.add_paragraph()
    nested_heading.paragraph_format.space_before = Pt(8)
    heading_run = nested_heading.add_run("NESTED TABLE CONTROL")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("1F5A92")

    outer = document.add_table(rows=1, cols=2)
    outer.style = "Table Grid"
    outer.cell(0, 0).text = "TABLE_NESTED_OUTER"
    nested_host = outer.cell(0, 1)
    nested_host.text = "TABLE_NESTED_HOST"
    nested = nested_host.add_table(rows=2, cols=2)
    nested.style = "Table Grid"
    nested_values = [["NESTED_A1", "NESTED_B1"], ["NESTED_A2", "NESTED_B2"]]
    for row_index, row in enumerate(nested.rows):
        for column_index, cell in enumerate(row.cells):
            cell.text = nested_values[row_index][column_index]
            cell.paragraphs[0].paragraph_format.space_after = Pt(0)
            cell.paragraphs[0].runs[0].font.size = Pt(8)

    control = document.add_paragraph()
    control.paragraph_format.space_before = Pt(8)
    control.add_run("PRESERVE_TABLES_UNRELATED_6F3C")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(TABLES_OUTPUT)
    with ZipFile(TABLES_OUTPUT, "a", compression=ZIP_DEFLATED, compresslevel=9) as package:
        package.writestr(
            "customXml/ctox-table-preserve.xml",
            b'<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="document.tables">TABLE_CUSTOM_PART_6F3C</ctox-preserve>',
        )
    canonicalize_zip(TABLES_OUTPUT)
    print(TABLES_OUTPUT)


def create_images_positioning() -> None:
    asset_directory = OUTPUT_DIRECTORY / ".generated-assets"
    asset_directory.mkdir(parents=True, exist_ok=True)
    inline_image = asset_directory / "inline-target.png"
    floating_image = asset_directory / "floating-target.png"
    preservation_image = asset_directory / "preserve-control.png"
    write_solid_png(inline_image, 240, 120, (31, 90, 146))
    write_solid_png(floating_image, 180, 180, (176, 80, 32))
    write_solid_png(preservation_image, 8, 8, (23, 107, 91))

    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10)
    normal.paragraph_format.space_after = Pt(5)
    document.core_properties.title = "CTOX Oracle Images Positioning Fixture"
    document.core_properties.subject = "document.images-positioning"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 11, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX IMAGES / POSITIONING ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("7A3E00")

    document.add_paragraph(
        "Resize the inline image and move the square-wrapped floating image while preserving media controls."
    )

    inline_heading = document.add_paragraph()
    inline_heading.add_run("INLINE_IMAGE_TARGET").bold = True
    inline_paragraph = document.add_paragraph()
    inline_paragraph.add_run("INLINE_BEFORE ")
    inline_shape = inline_paragraph.add_run().add_picture(str(inline_image), width=Inches(2.0))
    inline_shape._inline.docPr.set("descr", "CTOX_INLINE_IMAGE_TARGET")
    inline_paragraph.add_run(" INLINE_AFTER")

    floating_heading = document.add_paragraph()
    floating_heading.paragraph_format.space_before = Pt(14)
    floating_heading.add_run("FLOATING_IMAGE_TARGET").bold = True
    floating_paragraph = document.add_paragraph(
        "FLOATING_WRAP_TEXT_ALPHA FLOATING_WRAP_TEXT_BETA FLOATING_WRAP_TEXT_GAMMA "
        "FLOATING_WRAP_TEXT_DELTA FLOATING_WRAP_TEXT_EPSILON"
    )
    floating_shape = floating_paragraph.add_run().add_picture(str(floating_image), width=Inches(1.25))
    floating_shape._inline.docPr.set("descr", "CTOX_FLOATING_IMAGE_TARGET")
    make_floating_picture(floating_shape._inline, horizontal_twips=2200, vertical_twips=120)

    control = document.add_paragraph()
    control.paragraph_format.space_before = Pt(80)
    control.add_run("PRESERVE_IMAGES_UNRELATED_A91E")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(IMAGES_POSITIONING_OUTPUT)
    with ZipFile(IMAGES_POSITIONING_OUTPUT, "a", compression=ZIP_DEFLATED, compresslevel=9) as package:
        package.writestr("word/media/ctox-preserve-control.png", preservation_image.read_bytes())
        package.writestr(
            "customXml/ctox-image-preserve.xml",
            b'<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="document.images-positioning">IMAGE_CUSTOM_PART_A91E</ctox-preserve>',
        )
        package.writestr(
            "customXml/_rels/ctox-image-preserve.xml.rels",
            b'<?xml version="1.0" encoding="UTF-8"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdPreserveImage" Type="urn:ctox:preserve:image-control" Target="../../word/media/ctox-preserve-control.png"/></Relationships>',
        )
    canonicalize_zip(IMAGES_POSITIONING_OUTPUT)
    for asset in (inline_image, floating_image, preservation_image):
        asset.unlink()
    asset_directory.rmdir()
    print(IMAGES_POSITIONING_OUTPUT)


def create_sections_headers_footers() -> None:
    document = Document()
    section_one = document.sections[0]
    section_one.page_width = Inches(8.5)
    section_one.page_height = Inches(11)
    section_one.top_margin = Inches(0.8)
    section_one.bottom_margin = Inches(0.8)
    section_one.left_margin = Inches(0.9)
    section_one.right_margin = Inches(0.9)
    section_one.header_distance = Inches(0.4)
    section_one.footer_distance = Inches(0.45)
    section_one.different_first_page_header_footer = True

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(6)
    document.core_properties.title = "CTOX Oracle Sections Headers Footers Fixture"
    document.core_properties.subject = "document.sections-headers-footers"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 11, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    section_one.header.paragraphs[0].text = "HEADER_SECTION1_DEFAULT"
    section_one.header.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    section_one.first_page_header.paragraphs[0].text = "HEADER_SECTION1_FIRST"
    section_one.first_page_header.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    section_one.footer.paragraphs[0].text = "FOOTER_SECTION1_DEFAULT"
    section_one.footer.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX SECTIONS / HEADERS / FOOTERS ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("385723")
    document.add_paragraph(
        "Change the second section page setup and header/footer links, then insert a next-page section break."
    )
    document.add_paragraph("SECTION1_BODY_CONTROL")
    document.add_paragraph("PRESERVE_SECTION1_TEXT_38C1")

    section_two = document.add_section(WD_SECTION.NEW_PAGE)
    section_two.page_width = Inches(8.5)
    section_two.page_height = Inches(11)
    section_two.top_margin = Inches(0.8)
    section_two.bottom_margin = Inches(0.8)
    section_two.left_margin = Inches(0.9)
    section_two.right_margin = Inches(0.9)
    section_two.header_distance = Inches(0.4)
    section_two.footer_distance = Inches(0.45)
    section_two.header.is_linked_to_previous = True
    section_two.footer.is_linked_to_previous = True
    section_two.different_first_page_header_footer = False

    heading = document.add_paragraph()
    heading_run = heading.add_run("SECTION2_PAGE_SETUP_TARGET")
    heading_run.font.bold = True
    heading_run.font.color.rgb = RGBColor.from_string("385723")
    document.add_paragraph("SECTION2_HEADER_LINK_TARGET")
    document.add_paragraph("SECTION2_FOOTER_LINK_TARGET")
    document.add_paragraph("SECTION_BREAK_INSERT_AFTER_TARGET")
    document.add_paragraph("PRESERVE_SECTIONS_UNRELATED_4D72")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(SECTIONS_HEADERS_FOOTERS_OUTPUT)
    with ZipFile(SECTIONS_HEADERS_FOOTERS_OUTPUT, "a", compression=ZIP_DEFLATED, compresslevel=9) as package:
        package.writestr(
            "customXml/ctox-section-preserve.xml",
            b'<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="document.sections-headers-footers">SECTION_CUSTOM_PART_4D72</ctox-preserve>',
        )
    canonicalize_zip(SECTIONS_HEADERS_FOOTERS_OUTPUT)
    print(SECTIONS_HEADERS_FOOTERS_OUTPUT)


def create_links_bookmarks_fields() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(6)
    document.core_properties.title = "CTOX Oracle Links Bookmarks Fields Fixture"
    document.core_properties.subject = "document.links-bookmarks-fields"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 11, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX LINKS / BOOKMARKS / FIELDS ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("1F5A92")
    document.add_paragraph(
        "Create an external hyperlink and bookmark, then update the deterministic NUMPAGES field."
    )

    link_heading = document.add_paragraph()
    link_heading.add_run("HYPERLINK TARGET").bold = True
    link_target = document.add_paragraph("LINK_CREATE_TARGET")
    link_target.runs[0].font.size = Pt(14)
    link_control = document.add_paragraph("EXISTING_LINK_CONTROL: ")
    add_external_hyperlink(link_control, "CTOX_EXISTING_LINK", "https://ctox.dev/preserve-link")

    bookmark_heading = document.add_paragraph()
    bookmark_heading.paragraph_format.space_before = Pt(10)
    bookmark_heading.add_run("BOOKMARK TARGET").bold = True
    bookmark_target = document.add_paragraph("BOOKMARK_CREATE_TARGET")
    bookmark_target.runs[0].font.size = Pt(14)
    bookmark_control = document.add_paragraph("EXISTING_BOOKMARK_CONTROL: ")
    add_bookmark(bookmark_control, "ctox_existing_bookmark", 41, "CTOX_EXISTING_BOOKMARK")

    field_heading = document.add_paragraph()
    field_heading.paragraph_format.space_before = Pt(10)
    field_heading.add_run("FIELD TARGET").bold = True
    field_target = document.add_paragraph("NUMPAGES_FIELD_TARGET: ")
    add_complex_field(field_target, "NUMPAGES", "99")

    control = document.add_paragraph()
    control.paragraph_format.space_before = Pt(14)
    control.add_run("PRESERVE_LINKS_BOOKMARKS_FIELDS_UNRELATED_B73A")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(LINKS_BOOKMARKS_FIELDS_OUTPUT)
    with ZipFile(LINKS_BOOKMARKS_FIELDS_OUTPUT, "a", compression=ZIP_DEFLATED, compresslevel=9) as package:
        package.writestr(
            "customXml/ctox-links-preserve.xml",
            b'<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="document.links-bookmarks-fields">LINKS_CUSTOM_PART_B73A</ctox-preserve>',
        )
    canonicalize_zip(LINKS_BOOKMARKS_FIELDS_OUTPUT)
    print(LINKS_BOOKMARKS_FIELDS_OUTPUT)


def create_comments_track_changes() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(6)
    document.core_properties.title = "CTOX Oracle Comments Track Changes Fixture"
    document.core_properties.subject = "document.comments-track-changes"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 11, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX COMMENTS / TRACK CHANGES ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("7030A0")
    document.add_paragraph(
        "Create and resolve a comment, then record, accept, and reject deterministic revisions."
    )

    comment_heading = document.add_paragraph()
    comment_heading.add_run("COMMENT TARGETS").bold = True
    comment_target = document.add_paragraph("COMMENT_CREATE_TARGET")
    comment_target.runs[0].font.size = Pt(14)
    existing_comment = document.add_paragraph("EXISTING_COMMENT_CONTROL: ")
    add_comment_range(existing_comment, "CTOX_EXISTING_COMMENT_TEXT", 7)

    revision_heading = document.add_paragraph()
    revision_heading.paragraph_format.space_before = Pt(10)
    revision_heading.add_run("TRACK-CHANGE TARGETS").bold = True
    insert_target = document.add_paragraph("TRACK_INSERT_TARGET")
    insert_target.runs[0].font.size = Pt(14)
    delete_target = document.add_paragraph("TRACK_DELETE_TARGET")
    delete_target.runs[0].font.size = Pt(14)

    existing_revision = document.add_paragraph("EXISTING_REVISION_CONTROL: ")
    add_existing_revision(existing_revision, "ins", 51, "CTOX_EXISTING_INSERTION")
    existing_revision.add_run(" / ")
    add_existing_revision(existing_revision, "del", 52, "CTOX_EXISTING_DELETION")

    control = document.add_paragraph()
    control.paragraph_format.space_before = Pt(14)
    control.add_run("PRESERVE_COMMENTS_TRACK_CHANGES_UNRELATED_C9E4")

    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(COMMENTS_TRACK_CHANGES_OUTPUT)
    with ZipFile(COMMENTS_TRACK_CHANGES_OUTPUT, "r") as package:
        entries = {item.filename: package.read(item.filename) for item in package.infolist()}
    relationship = (
        b'<Relationship Id="rIdCtoxComments" '
        b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" '
        b'Target="comments.xml"/>'
    )
    assert b"</Relationships>" in entries["word/_rels/document.xml.rels"]
    entries["word/_rels/document.xml.rels"] = entries["word/_rels/document.xml.rels"].replace(
        b"</Relationships>", relationship + b"</Relationships>"
    )
    override = (
        b'<Override PartName="/word/comments.xml" '
        b'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.comments+xml"/>'
    )
    assert b"</Types>" in entries["[Content_Types].xml"]
    entries["[Content_Types].xml"] = entries["[Content_Types].xml"].replace(
        b"</Types>", override + b"</Types>"
    )
    entries["word/comments.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        b'<w:comment w:id="7" w:author="CTOX Existing Reviewer" w:initials="CE" '
        b'w:date="2026-07-11T00:00:00Z"><w:p><w:r><w:t>'
        b'CTOX_EXISTING_COMMENT_BODY</w:t></w:r></w:p></w:comment></w:comments>'
    )
    entries["customXml/ctox-comments-preserve.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8"?><ctox-preserve '
        b'feature="document.comments-track-changes">COMMENTS_CUSTOM_PART_C9E4</ctox-preserve>'
    )
    temporary = COMMENTS_TRACK_CHANGES_OUTPUT.with_suffix(".parts.docx")
    with ZipFile(temporary, "w", compression=ZIP_DEFLATED, compresslevel=9) as package:
        for name, data in sorted(entries.items()):
            package.writestr(name, data)
    temporary.replace(COMMENTS_TRACK_CHANGES_OUTPUT)
    canonicalize_zip(COMMENTS_TRACK_CHANGES_OUTPUT)
    print(COMMENTS_TRACK_CHANGES_OUTPUT)


def create_drawings_charts() -> None:
    """Create a DOCX with a real chart, a VML business shape, and escrow controls."""
    asset_directory = OUTPUT_DIRECTORY / ".generated-drawings-charts-assets"
    asset_directory.mkdir(parents=True, exist_ok=True)
    chart_source = asset_directory / "chart-source.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart_data = ChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Revenue", (12, 18, 27))
    chart_data.add_series("Costs", (8, 11, 15))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        PptxInches(1),
        PptxInches(1),
        PptxInches(6),
        PptxInches(3.4),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "CTOX Quarterly Business"
    chart.has_legend = True
    presentation.save(chart_source)

    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.top_margin = Inches(0.6)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10)
    normal.paragraph_format.space_after = Pt(5)
    document.core_properties.title = "CTOX Oracle Drawings Charts Fixture"
    document.core_properties.subject = "document.drawings-charts"
    document.core_properties.author = "CTOX"
    document.core_properties.last_modified_by = "CTOX"
    fixed_timestamp = datetime(2026, 7, 11, 0, 0, 0, tzinfo=timezone.utc)
    document.core_properties.created = fixed_timestamp
    document.core_properties.modified = fixed_timestamp

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CTOX DRAWINGS / CHARTS ORACLE")
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("5B2C83")
    document.add_paragraph(
        "Resize and restyle the business shape, then edit and reposition the chart while preserving unrelated parts."
    )

    shape_heading = document.add_paragraph()
    shape_heading.add_run("DRAWING_SHAPE_TARGET").bold = True
    shape_paragraph = document.add_paragraph()
    shape_run = etree.fromstring(b"""
      <w:r xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
           xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
           xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">
        <w:drawing>
          <wp:inline distT="0" distB="0" distL="0" distR="0">
            <wp:extent cx="2286000" cy="685800"/>
            <wp:effectExtent l="6350" t="6350" r="6350" b="6350"/>
            <wp:docPr id="90" name="CTOX_EXISTING_BUSINESS_SHAPE"/>
            <wp:cNvGraphicFramePr/>
            <a:graphic>
              <a:graphicData uri="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">
                <wps:wsp>
                  <wps:cNvPr id="0" name="CTOX_EXISTING_BUSINESS_SHAPE"/>
                  <wps:cNvSpPr/>
                  <wps:spPr>
                    <a:xfrm rot="0"><a:off x="0" y="0"/><a:ext cx="2286000" cy="685800"/></a:xfrm>
                    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
                    <a:solidFill><a:srgbClr val="D9EAD3"/></a:solidFill>
                    <a:ln w="12700"><a:solidFill><a:srgbClr val="176B5B"/></a:solidFill></a:ln>
                  </wps:spPr>
                  <wps:txbx><w:txbxContent><w:p><w:r><w:t>CTOX_EXISTING_BUSINESS_SHAPE</w:t></w:r></w:p></w:txbxContent></wps:txbx>
                  <wps:bodyPr wrap="square" lIns="12700" tIns="12700" rIns="12700" bIns="12700" upright="1"/>
                </wps:wsp>
              </a:graphicData>
            </a:graphic>
          </wp:inline>
        </w:drawing>
      </w:r>
    """)
    shape_paragraph._p.append(shape_run)

    chart_heading = document.add_paragraph()
    chart_heading.paragraph_format.space_before = Pt(10)
    chart_heading.add_run("CHART_OBJECT_TARGET").bold = True
    chart_paragraph = document.add_paragraph()
    chart_run = OxmlElement("w:r")
    drawing = OxmlElement("w:drawing")
    inline = OxmlElement("wp:inline")
    for name in ("distT", "distB", "distL", "distR"):
        inline.set(name, "0")
    extent = OxmlElement("wp:extent")
    extent.set("cx", "5486400")
    extent.set("cy", "2743200")
    effect = OxmlElement("wp:effectExtent")
    for name in ("l", "t", "r", "b"):
        effect.set(name, "0")
    doc_properties = OxmlElement("wp:docPr")
    doc_properties.set("id", "91")
    doc_properties.set("name", "CTOX_CHART_TARGET")
    frame_properties = OxmlElement("wp:cNvGraphicFramePr")
    locks = OxmlElement("a:graphicFrameLocks")
    locks.set("noChangeAspect", "1")
    frame_properties.append(locks)
    graphic = OxmlElement("a:graphic")
    graphic_data = OxmlElement("a:graphicData")
    graphic_data.set("uri", "http://schemas.openxmlformats.org/drawingml/2006/chart")
    chart_reference = OxmlElement("c:chart")
    chart_reference.set(qn("r:id"), "rIdCtoxChart")
    graphic_data.append(chart_reference)
    graphic.append(graphic_data)
    for child in (extent, effect, doc_properties, frame_properties, graphic):
        inline.append(child)
    drawing.append(inline)
    chart_run.append(drawing)
    chart_paragraph._p.append(chart_run)

    control = document.add_paragraph()
    control.add_run("PRESERVE_DRAWINGS_CHARTS_UNRELATED_D4A7")
    OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    document.save(DRAWINGS_CHARTS_OUTPUT)

    with ZipFile(chart_source, "r") as source:
        chart_xml = source.read("ppt/charts/chart1.xml")
        chart_rels = source.read("ppt/charts/_rels/chart1.xml.rels")
        workbook_name = next(
            item.filename for item in source.infolist()
            if item.filename.startswith("ppt/embeddings/") and item.filename.endswith(".xlsx")
        )
        workbook = source.read(workbook_name)
    # XlsxWriter stamps the embedded chart workbook with wall-clock time. Pin
    # that one volatile value and preserve its deterministic entry order and
    # metadata so repeated fixture builds are byte-identical.
    with ZipFile(BytesIO(workbook), "r") as source:
        workbook_entries = [(item, source.read(item.filename)) for item in source.infolist()]
    normalized_workbook = BytesIO()
    with ZipFile(normalized_workbook, "w", compression=ZIP_DEFLATED) as target:
        for item, data in workbook_entries:
            if item.filename == "docProps/core.xml":
                data = re.sub(
                    rb"\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}Z",
                    b"2026-07-11T01:19:41Z",
                    data,
                )
            target.writestr(item, data)
    workbook = normalized_workbook.getvalue()

    with ZipFile(DRAWINGS_CHARTS_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    entries["word/charts/chart1.xml"] = chart_xml
    entries["word/charts/_rels/chart1.xml.rels"] = chart_rels
    entries["word/embeddings/ctox-chart-data.xlsx"] = workbook
    entries["customXml/ctox-drawings-charts-preserve.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="document.drawings-charts">DRAWINGS_CHARTS_CUSTOM_PART_D4A7</ctox-preserve>'
    )
    entries["word/charts/_rels/chart1.xml.rels"] = entries[
        "word/charts/_rels/chart1.xml.rels"
    ].replace(
        workbook_name.removeprefix("ppt/embeddings/").encode(),
        b"ctox-chart-data.xlsx",
    )
    content_types = entries["[Content_Types].xml"].replace(
        b"</Types>",
        b'<Override PartName="/word/charts/chart1.xml" ContentType="application/vnd.openxmlformats-officedocument.drawingml.chart+xml"/>'
        b'<Override PartName="/word/embeddings/ctox-chart-data.xlsx" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"/>'
        b"</Types>",
    )
    entries["[Content_Types].xml"] = content_types
    document_rels = entries["word/_rels/document.xml.rels"].replace(
        b"</Relationships>",
        b'<Relationship Id="rIdCtoxChart" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" Target="charts/chart1.xml"/>'
        b"</Relationships>",
    )
    entries["word/_rels/document.xml.rels"] = document_rels
    temporary = DRAWINGS_CHARTS_OUTPUT.with_suffix(".assembled.docx")
    with ZipFile(temporary, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    temporary.replace(DRAWINGS_CHARTS_OUTPUT)
    canonicalize_zip(DRAWINGS_CHARTS_OUTPUT)
    chart_source.unlink()
    asset_directory.rmdir()
    print(DRAWINGS_CHARTS_OUTPUT)


def create_spreadsheet_open_render_sheets() -> None:
    """Create a deterministic three-sheet XLSX without a generator runtime dependency."""
    SPREADSHEET_OUTPUT_DIRECTORY.mkdir(parents=True, exist_ok=True)
    shared_values = [
        "CTOX OVERVIEW", "ORACLE_SHEET_OVERVIEW_71C9", "Region", "Revenue", "North", "South",
        "CTOX DETAILS", "ORACLE_SHEET_DETAILS_24AF", "Item", "Quantity", "Consulting", "Support",
        "CTOX ARCHIVE", "ORACLE_SHEET_ARCHIVE_B83D", "Year", "Status", "2025", "Closed",
    ]
    shared_strings = "".join(f"<si><t>{value}</t></si>" for value in shared_values)
    worksheets = [
        ("Overview", "0", "1", [("A3", 2, 0), ("B3", 3, 0), ("A4", 4, 0), ("B4", 125000, 1), ("A5", 5, 0), ("B5", 98000, 1)]),
        ("Details", "6", "7", [("A3", 8, 0), ("B3", 9, 0), ("A4", 10, 0), ("B4", 42, 1), ("A5", 11, 0), ("B5", 18, 1)]),
        ("Archive", "12", "13", [("A3", 14, 0), ("B3", 15, 0), ("A4", 16, 0), ("B4", 17, 0)]),
    ]
    entries: dict[str, bytes] = {
        "[Content_Types].xml": b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/><Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/><Override PartName="/xl/sharedStrings.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sharedStrings+xml"/><Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/><Override PartName="/xl/worksheets/sheet2.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/><Override PartName="/xl/worksheets/sheet3.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/><Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/><Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/></Types>''',
        "_rels/.rels": b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/><Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/></Relationships>''',
        "docProps/core.xml": b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"><dc:title>CTOX Oracle Spreadsheet Open Render Sheets</dc:title><dc:creator>CTOX</dc:creator><dcterms:created xsi:type="dcterms:W3CDTF">2026-07-11T00:00:00Z</dcterms:created><dcterms:modified xsi:type="dcterms:W3CDTF">2026-07-11T00:00:00Z</dcterms:modified></cp:coreProperties>''',
        "docProps/app.xml": b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"><Application>CTOX Oracle Fixture</Application><TitlesOfParts><vt:vector size="3" baseType="lpstr"><vt:lpstr>Overview</vt:lpstr><vt:lpstr>Details</vt:lpstr><vt:lpstr>Archive</vt:lpstr></vt:vector></TitlesOfParts></Properties>''',
        "xl/workbook.xml": b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><bookViews><workbookView activeTab="0"/></bookViews><sheets><sheet name="Overview" sheetId="1" r:id="rId1"/><sheet name="Details" sheetId="2" r:id="rId2"/><sheet name="Archive" sheetId="3" state="hidden" r:id="rId3"/></sheets><calcPr calcId="191029"/></workbook>''',
        "xl/_rels/workbook.xml.rels": b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet2.xml"/><Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet3.xml"/><Relationship Id="rId4" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/><Relationship Id="rId5" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/sharedStrings" Target="sharedStrings.xml"/></Relationships>''',
        "xl/sharedStrings.xml": f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="{len(shared_values)}" uniqueCount="{len(shared_values)}">{shared_strings}</sst>'''.encode(),
        "xl/styles.xml": b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><fonts count="2"><font><sz val="11"/><name val="Aptos"/></font><font><b/><color rgb="FFFFFFFF"/><sz val="16"/><name val="Aptos Display"/></font></fonts><fills count="3"><fill><patternFill patternType="none"/></fill><fill><patternFill patternType="gray125"/></fill><fill><patternFill patternType="solid"><fgColor rgb="FF176B5B"/><bgColor indexed="64"/></patternFill></fill></fills><borders count="1"><border><left/><right/><top/><bottom/><diagonal/></border></borders><cellStyleXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0"/></cellStyleXfs><cellXfs count="3"><xf numFmtId="0" fontId="0" fillId="0" borderId="0" xfId="0"/><xf numFmtId="0" fontId="0" fillId="0" borderId="0" xfId="0"/><xf numFmtId="0" fontId="1" fillId="2" borderId="0" xfId="0" applyFont="1" applyFill="1" applyAlignment="1"><alignment horizontal="center"/></xf></cellXfs><cellStyles count="1"><cellStyle name="Normal" xfId="0" builtinId="0"/></cellStyles></styleSheet>''',
        "customXml/ctox-spreadsheet-preserve.xml": b'''<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="spreadsheet.open-render-sheets">SPREADSHEET_CUSTOM_PART_5E91</ctox-preserve>''',
    }
    for index, (_, title_index, marker_index, cells) in enumerate(worksheets, 1):
        rows = [f'<row r="1" ht="26" customHeight="1"><c r="A1" t="s" s="2"><v>{title_index}</v></c></row>', f'<row r="2"><c r="A2" t="s"><v>{marker_index}</v></c></row>']
        by_row: dict[int, list[str]] = {}
        for coordinate, value, numeric in cells:
            row = int("".join(filter(str.isdigit, coordinate)))
            cell = f'<c r="{coordinate}"><v>{value}</v></c>' if numeric else f'<c r="{coordinate}" t="s"><v>{value}</v></c>'
            by_row.setdefault(row, []).append(cell)
        rows.extend(f'<row r="{row}">{"".join(values)}</row>' for row, values in sorted(by_row.items()))
        entries[f"xl/worksheets/sheet{index}.xml"] = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><dimension ref="A1:B5"/><sheetViews><sheetView workbookViewId="0"><selection activeCell="A1" sqref="A1"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="1" width="24" customWidth="1"/><col min="2" max="2" width="18" customWidth="1"/></cols><sheetData>{"".join(rows)}</sheetData></worksheet>'''.encode()
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_OPEN_RENDER_OUTPUT)


def create_spreadsheet_edit_save() -> None:
    """Derive a deterministic editable XLSX with target and escrow markers."""
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    replacements = {
        b"CTOX OVERVIEW": b"CTOX EDIT SAVE",
        b"ORACLE_SHEET_OVERVIEW_71C9": b"CTOX_EDIT_CELL_ALPHA",
        b"Region": b"PRESERVE_XLSX_UNRELATED_6D2A",
        b"CTOX Oracle Spreadsheet Open Render Sheets": b"CTOX Oracle Spreadsheet Edit Save",
        b"spreadsheet.open-render-sheets": b"spreadsheet.edit-save",
        b"SPREADSHEET_CUSTOM_PART_5E91": b"SPREADSHEET_EDIT_ESCROW_81F4",
    }
    for path in ("xl/sharedStrings.xml", "docProps/core.xml", "customXml/ctox-spreadsheet-preserve.xml"):
        for before, after in replacements.items():
            entries[path] = entries[path].replace(before, after)
    with ZipFile(SPREADSHEET_EDIT_SAVE_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_EDIT_SAVE_OUTPUT)


def create_spreadsheet_undo_clipboard_fill() -> None:
    """Create deterministic cells for history, clipboard and range-fill flows."""
    with ZipFile(SPREADSHEET_EDIT_SAVE_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    replacements = {
        b"CTOX EDIT SAVE": b"CTOX UNDO CLIPBOARD FILL",
        b"CTOX_EDIT_CELL_ALPHA": b"UNDO_FILL_BASE",
        b"PRESERVE_XLSX_UNRELATED_6D2A": b"COPY_SOURCE_TEXT",
        b"Revenue": b"PASTE_TARGET_TEXT",
        b"North": b"FILL_ROW_ONE",
        b"South": b"FILL_ROW_TWO",
        b"Spreadsheet Edit Save": b"Spreadsheet Undo Clipboard Fill",
        b"spreadsheet.edit-save": b"spreadsheet.undo-clipboard-fill",
        b"SPREADSHEET_EDIT_ESCROW_81F4": b"SPREADSHEET_UNDO_ESCROW_9A31",
    }
    for path in ("xl/sharedStrings.xml", "docProps/core.xml", "customXml/ctox-spreadsheet-preserve.xml"):
        for before, after in replacements.items():
            entries[path] = entries[path].replace(before, after)
    with ZipFile(SPREADSHEET_UNDO_CLIPBOARD_FILL_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_UNDO_CLIPBOARD_FILL_OUTPUT)


def create_spreadsheet_cell_format_rows_columns() -> None:
    """Create deterministic targets for cell styles and row/column geometry."""
    with ZipFile(SPREADSHEET_UNDO_CLIPBOARD_FILL_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    replacements = {
        b"CTOX UNDO CLIPBOARD FILL": b"CTOX CELL FORMAT ROWS COLUMNS",
        b"UNDO_FILL_BASE": b"FORMAT_TEXT_TARGET",
        b"COPY_SOURCE_TEXT": b"CURRENCY_VALUE_TARGET",
        b"PASTE_TARGET_TEXT": b"COLUMN_WIDTH_TARGET",
        b"FILL_ROW_ONE": b"ROW_HEIGHT_TARGET",
        b"FILL_ROW_TWO": b"VISIBILITY_TARGET",
        b"Spreadsheet Undo Clipboard Fill": b"Spreadsheet Cell Format Rows Columns",
        b"spreadsheet.undo-clipboard-fill": b"spreadsheet.cell-format-rows-columns",
        b"SPREADSHEET_UNDO_ESCROW_9A31": b"SPREADSHEET_FORMAT_ESCROW_4C72",
    }
    for path in ("xl/sharedStrings.xml", "docProps/core.xml", "customXml/ctox-spreadsheet-preserve.xml"):
        for before, after in replacements.items():
            entries[path] = entries[path].replace(before, after)
    with ZipFile(SPREADSHEET_CELL_FORMAT_ROWS_COLUMNS_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_CELL_FORMAT_ROWS_COLUMNS_OUTPUT)


def create_spreadsheet_formulas_references() -> None:
    """Create deterministic formula, reference-shift and error-value targets."""
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    shared_values = [
        "CTOX FORMULAS REFERENCES", "INPUT_BASE", "RELATIVE_FORMULA",
        "ABSOLUTE_FORMULA", "RANGE_FORMULA", "SHEET_REFERENCE",
        "CTOX DETAILS", "ORACLE_SHEET_DETAILS_24AF", "Item", "Quantity",
        "Consulting", "Support", "CTOX ARCHIVE", "ORACLE_SHEET_ARCHIVE_B83D",
        "Year", "Status", "2025", "Closed", "COPY_SOURCE", "ERROR_FORMULA",
        "MANUAL_FORMULA",
    ]
    shared_strings = "".join(f"<si><t>{value}</t></si>" for value in shared_values)
    entries["xl/sharedStrings.xml"] = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="{len(shared_values)}" uniqueCount="{len(shared_values)}">{shared_strings}</sst>'''.encode()
    entries["xl/worksheets/sheet1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><dimension ref="A1:D8"/><sheetViews><sheetView workbookViewId="0"><selection activeCell="A1" sqref="A1"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="1" width="24" customWidth="1"/><col min="2" max="4" width="18" customWidth="1"/></cols><sheetData><row r="1" ht="26" customHeight="1"><c r="A1" t="s" s="2"><v>0</v></c></row><row r="2"><c r="A2" t="s"><v>1</v></c><c r="B2"><v>10</v></c><c r="C2"><v>20</v></c><c r="D2"><v>5</v></c></row><row r="3"><c r="A3" t="s"><v>2</v></c><c r="B3"><f>B2*2</f><v>20</v></c><c r="D3"><f>D2*2</f><v>10</v></c></row><row r="4"><c r="A4" t="s"><v>3</v></c><c r="B4"><f>$B$2+5</f><v>15</v></c></row><row r="5"><c r="A5" t="s"><v>4</v></c><c r="B5"><f>SUM(B2:B4)</f><v>45</v></c></row><row r="6"><c r="A6" t="s"><v>5</v></c><c r="B6"><f>'Details'!B4+1</f><v>43</v></c></row><row r="7"><c r="A7" t="s"><v>18</v></c><c r="B7"><f>B2+1</f><v>11</v></c><c r="C7"><f>B2+1</f><v>11</v></c></row><row r="8"><c r="A8" t="s"><v>19</v></c><c r="B8" t="e"><f>1/0</f><v>#DIV/0!</v></c><c r="D8" t="s"><v>20</v></c></row></sheetData></worksheet>'''
    entries["docProps/core.xml"] = entries["docProps/core.xml"].replace(
        b"Spreadsheet Open Render Sheets", b"Spreadsheet Formulas References"
    )
    entries["customXml/ctox-spreadsheet-preserve.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="spreadsheet.formulas-references">SPREADSHEET_FORMULA_ESCROW_73D1</ctox-preserve>'
    )
    with ZipFile(SPREADSHEET_FORMULAS_REFERENCES_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_FORMULAS_REFERENCES_OUTPUT)


def create_spreadsheet_multi_sheet_merge_freeze() -> None:
    """Create deterministic multi-sheet, merged-cell and frozen-pane targets."""
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    shared_values = [
        "CTOX MULTI SHEET MERGE FREEZE", "OVERVIEW_MARKER_6F21", "MERGED_SOURCE",
        "MERGE_TARGET", "FREEZE_ANCHOR", "CTOX OPERATIONS", "OPERATIONS_MARKER_A9C4",
        "SECOND_SHEET_CELL", "CTOX ARCHIVE", "ARCHIVE_HIDDEN_MARKER_13D8",
    ]
    strings = "".join(f"<si><t>{value}</t></si>" for value in shared_values)
    entries["xl/sharedStrings.xml"] = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="{len(shared_values)}" uniqueCount="{len(shared_values)}">{strings}</sst>'''.encode()
    entries["xl/worksheets/sheet1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><dimension ref="A1:D7"/><sheetViews><sheetView workbookViewId="0"><pane xSplit="1" ySplit="1" topLeftCell="B2" activePane="bottomRight" state="frozen"/><selection pane="bottomRight" activeCell="B2" sqref="B2"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="4" width="20" customWidth="1"/></cols><sheetData><row r="1" ht="26" customHeight="1"><c r="A1" t="s" s="2"><v>0</v></c></row><row r="2"><c r="A2" t="s"><v>1</v></c><c r="B2" t="s"><v>2</v></c></row><row r="3"><c r="A3" t="s"><v>3</v></c><c r="B3"/><c r="C3"/></row><row r="4"><c r="A4" t="s"><v>4</v></c><c r="B4"><v>42</v></c></row><row r="7"><c r="D7"><v>73</v></c></row></sheetData><mergeCells count="1"><mergeCell ref="B2:C2"/></mergeCells></worksheet>'''
    entries["xl/worksheets/sheet2.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><dimension ref="A1:C5"/><sheetViews><sheetView workbookViewId="0"><selection activeCell="A1" sqref="A1"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="3" width="20" customWidth="1"/></cols><sheetData><row r="1"><c r="A1" t="s" s="2"><v>5</v></c></row><row r="2"><c r="A2" t="s"><v>6</v></c></row><row r="4"><c r="B4" t="s"><v>7</v></c></row><row r="5"><c r="C5"><v>19</v></c></row></sheetData></worksheet>'''
    entries["xl/worksheets/sheet3.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><dimension ref="A1:B2"/><sheetViews><sheetView workbookViewId="0"><selection activeCell="A1" sqref="A1"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><sheetData><row r="1"><c r="A1" t="s" s="2"><v>8</v></c></row><row r="2"><c r="A2" t="s"><v>9</v></c></row></sheetData></worksheet>'''
    entries["docProps/core.xml"] = entries["docProps/core.xml"].replace(b"Spreadsheet Open Render Sheets", b"Spreadsheet Multi Sheet Merge Freeze")
    entries["customXml/ctox-spreadsheet-preserve.xml"] = b'''<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="spreadsheet.multi-sheet-merge-freeze">SPREADSHEET_MERGE_FREEZE_ESCROW_D52B</ctox-preserve>'''
    with ZipFile(SPREADSHEET_MULTI_SHEET_MERGE_FREEZE_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0)); info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_MULTI_SHEET_MERGE_FREEZE_OUTPUT)


def create_spreadsheet_sort_filter_tables() -> None:
    """Create a deterministic structured table with sortable/filterable rows."""
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    # sheet2/sheet3 are inherited from the open-render fixture and reference
    # indices through 17. Keep the complete shared-string domain valid instead
    # of relying on DocumentServer's formerly permissive out-of-range fallback.
    shared_values = [
        "Region", "Product", "Revenue", "South", "Support", "North",
        "Consulting", "East", "Training", "West", "Audit", "Archive",
        "CTOX ARCHIVE", "ARCHIVE_TABLE_MARKER_8B47", "Year", "Status",
        "2025", "Closed",
    ]
    entries["xl/sharedStrings.xml"] = (f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="{len(shared_values)}" uniqueCount="{len(shared_values)}">''' + "".join(f"<si><t>{v}</t></si>" for v in shared_values) + "</sst>").encode()
    entries["xl/worksheets/sheet1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><dimension ref="A1:C6"/><sheetViews><sheetView workbookViewId="0"><selection activeCell="A1" sqref="A1:C6"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="2" width="18" customWidth="1"/><col min="3" max="3" width="14" customWidth="1"/></cols><sheetData><row r="1"><c r="A1" t="s" s="2"><v>0</v></c><c r="B1" t="s" s="2"><v>1</v></c><c r="C1" t="s" s="2"><v>2</v></c></row><row r="2"><c r="A2" t="s"><v>3</v></c><c r="B2" t="s"><v>4</v></c><c r="C2"><v>120</v></c></row><row r="3"><c r="A3" t="s"><v>5</v></c><c r="B3" t="s"><v>6</v></c><c r="C3"><v>420</v></c></row><row r="4"><c r="A4" t="s"><v>7</v></c><c r="B4" t="s"><v>8</v></c><c r="C4"><v>240</v></c></row><row r="5"><c r="A5" t="s"><v>5</v></c><c r="B5" t="s"><v>4</v></c><c r="C5"><v>310</v></c></row><row r="6"><c r="A6" t="s"><v>9</v></c><c r="B6" t="s"><v>10</v></c><c r="C6"><v>180</v></c></row></sheetData><tableParts count="1"><tablePart r:id="rId1"/></tableParts></worksheet>'''
    entries["xl/worksheets/_rels/sheet1.xml.rels"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/table" Target="../tables/table1.xml"/></Relationships>'''
    entries["xl/tables/table1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><table xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" id="1" name="RevenueTable" displayName="RevenueTable" ref="A1:C6" totalsRowShown="0"><autoFilter ref="A1:C6"/><tableColumns count="3"><tableColumn id="1" name="Region"/><tableColumn id="2" name="Product"/><tableColumn id="3" name="Revenue"/></tableColumns><tableStyleInfo name="TableStyleMedium4" showFirstColumn="0" showLastColumn="0" showRowStripes="1" showColumnStripes="0"/></table>'''
    entries["[Content_Types].xml"] = entries["[Content_Types].xml"].replace(b"</Types>", b'<Override PartName="/xl/tables/table1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.table+xml"/></Types>')
    entries["docProps/core.xml"] = entries["docProps/core.xml"].replace(b"Spreadsheet Open Render Sheets", b"Spreadsheet Sort Filter Tables")
    entries["customXml/ctox-spreadsheet-preserve.xml"] = b'''<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="spreadsheet.sort-filter-tables">SPREADSHEET_TABLE_ESCROW_8B47</ctox-preserve>'''
    with ZipFile(SPREADSHEET_SORT_FILTER_TABLES_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0)); info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_SORT_FILTER_TABLES_OUTPUT)


def create_spreadsheet_validation_conditional_formatting() -> None:
    """Create deterministic list/whole-number validation and conditional formats."""
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    # Preserve the inherited Details/Archive shared-string indices 6..17 and
    # allocate the feature-specific A6 marker at 18.
    values = [
        "CTOX VALIDATION CONDITIONAL", "Status", "Quantity", "Score Scale",
        "Threshold", "Draft", "CTOX DETAILS", "ORACLE_SHEET_DETAILS_24AF",
        "Item", "Quantity", "Consulting", "Support", "CTOX ARCHIVE",
        "ORACLE_SHEET_ARCHIVE_B83D", "Year", "Status", "2025", "Closed",
        "INPUT_TARGET",
    ]
    entries["xl/sharedStrings.xml"] = (f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="{len(values)}" uniqueCount="{len(values)}">''' + "".join(f"<si><t>{v}</t></si>" for v in values) + "</sst>").encode()
    entries["xl/worksheets/sheet1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><dimension ref="A1:E6"/><sheetViews><sheetView workbookViewId="0"><selection activeCell="B2" sqref="B2"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="1" width="26" customWidth="1"/><col min="2" max="5" width="16" customWidth="1"/></cols><sheetData><row r="1"><c r="A1" t="s" s="2"><v>0</v></c></row><row r="2"><c r="A2" t="s"><v>1</v></c><c r="B2" t="s"><v>5</v></c><c r="C2"><v>5</v></c><c r="D2"><v>10</v></c><c r="E2"><v>20</v></c></row><row r="3"><c r="A3" t="s"><v>2</v></c><c r="D3"><v>35</v></c><c r="E3"><v>55</v></c></row><row r="4"><c r="A4" t="s"><v>3</v></c><c r="D4"><v>60</v></c><c r="E4"><v>75</v></c></row><row r="5"><c r="A5" t="s"><v>4</v></c><c r="D5"><v>85</v></c><c r="E5"><v>40</v></c></row><row r="6"><c r="A6" t="s"><v>18</v></c><c r="D6"><v>100</v></c><c r="E6"><v>90</v></c></row></sheetData><conditionalFormatting sqref="D2:D6"><cfRule type="colorScale" priority="1"><colorScale><cfvo type="min"/><cfvo type="percentile" val="50"/><cfvo type="max"/><color rgb="FFF8696B"/><color rgb="FFFFEB84"/><color rgb="FF63BE7B"/></colorScale></cfRule></conditionalFormatting><conditionalFormatting sqref="E2:E6"><cfRule type="cellIs" dxfId="0" priority="2" operator="greaterThan"><formula>50</formula></cfRule></conditionalFormatting><dataValidations count="2"><dataValidation type="list" allowBlank="0" showErrorMessage="1" errorStyle="stop" errorTitle="Invalid status" error="Choose Draft, Review or Final" sqref="B2"><formula1>"Draft,Review,Final"</formula1></dataValidation><dataValidation type="whole" operator="between" allowBlank="0" showErrorMessage="1" errorStyle="stop" errorTitle="Invalid quantity" error="Enter a whole number from 1 to 10" sqref="C2"><formula1>1</formula1><formula2>10</formula2></dataValidation></dataValidations></worksheet>'''
    entries["xl/styles.xml"] = entries["xl/styles.xml"].replace(b"</styleSheet>", b'<dxfs count="1"><dxf><fill><patternFill patternType="solid"><fgColor rgb="FFC6EFCE"/><bgColor rgb="FFC6EFCE"/></patternFill></fill><font><color rgb="FF006100"/></font></dxf></dxfs></styleSheet>')
    entries["docProps/core.xml"] = entries["docProps/core.xml"].replace(b"Spreadsheet Open Render Sheets", b"Spreadsheet Validation Conditional Formatting")
    entries["customXml/ctox-spreadsheet-preserve.xml"] = b'''<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="spreadsheet.validation-conditional-formatting">SPREADSHEET_VALIDATION_ESCROW_C19E</ctox-preserve>'''
    with ZipFile(SPREADSHEET_VALIDATION_CONDITIONAL_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 11, 0, 0, 0)); info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_VALIDATION_CONDITIONAL_OUTPUT)


def create_spreadsheet_comments_names_protection() -> None:
    """Create deterministic classic-comment, defined-name and protection targets."""
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    # The inherited Details/Archive sheets reference shared-string indices up
    # through 17, so keep that complete domain and append feature markers.
    values = [
        "CTOX COMMENTS NAMES PROTECTION", "Named amount", "Protected input",
        "Comment target", "RESERVED_4", "RESERVED_5", "CTOX DETAILS", "ORACLE_SHEET_DETAILS_24AF", "Item",
        "Quantity", "Consulting", "Support", "CTOX ARCHIVE",
        "ORACLE_SHEET_ARCHIVE_B83D", "Year", "Status", "2025", "Closed",
        "COMMENT_CELL_MARKER",
    ]
    entries["xl/sharedStrings.xml"] = (f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="{len(values)}" uniqueCount="{len(values)}">''' + "".join(f"<si><t>{v}</t></si>" for v in values) + "</sst>").encode()
    entries["xl/worksheets/sheet1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><dimension ref="A1:C5"/><sheetViews><sheetView workbookViewId="0"><selection activeCell="B4" sqref="B4"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="1" width="28" customWidth="1"/><col min="2" max="3" width="18" customWidth="1"/></cols><sheetData><row r="1"><c r="A1" t="s" s="2"><v>0</v></c></row><row r="2"><c r="A2" t="s"><v>1</v></c><c r="B2"><v>1250</v></c></row><row r="3"><c r="A3" t="s"><v>2</v></c><c r="B3"><v>42</v></c></row><row r="4"><c r="A4" t="s"><v>3</v></c><c r="B4" t="s"><v>18</v></c></row></sheetData><sheetProtection sheet="1" objects="1" scenarios="1" formatCells="1" formatColumns="1" formatRows="1" insertColumns="1" insertRows="1" insertHyperlinks="1" deleteColumns="1" deleteRows="1" selectLockedCells="0" sort="1" autoFilter="1" pivotTables="1" selectUnlockedCells="0"/><legacyDrawing r:id="rId2"/></worksheet>'''
    entries["xl/worksheets/_rels/sheet1.xml.rels"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments" Target="../comments1.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/vmlDrawing" Target="../drawings/vmlDrawing1.vml"/></Relationships>'''
    entries["xl/comments1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><comments xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><authors><author>CTOX</author></authors><commentList><comment ref="B4" authorId="0"><text><r><rPr><b/><sz val="10"/><rFont val="Arial"/></rPr><t>CTOX_EXISTING_CELL_COMMENT</t></r></text></comment></commentList></comments>'''
    entries["xl/drawings/vmlDrawing1.vml"] = b'''<xml xmlns:v="urn:schemas-microsoft-com:vml" xmlns:o="urn:schemas-microsoft-com:office:office" xmlns:x="urn:schemas-microsoft-com:office:excel"><o:shapelayout v:ext="edit"><o:idmap v:ext="edit" data="1"/></o:shapelayout><v:shapetype id="_x0000_t202" coordsize="21600,21600" o:spt="202" path="m,l,21600r21600,l21600,xe"><v:stroke joinstyle="miter"/><v:path gradientshapeok="t" o:connecttype="rect"/></v:shapetype><v:shape id="_x0000_s1025" type="#_x0000_t202" style="position:absolute;margin-left:80pt;margin-top:5pt;width:108pt;height:59pt;z-index:1;visibility:hidden" fillcolor="#ffffe1" o:insetmode="auto"><v:fill color2="#ffffe1"/><v:shadow on="t" color="black" obscured="t"/><v:path o:connecttype="none"/><v:textbox style="mso-direction-alt:auto"><div style="text-align:left"/></v:textbox><x:ClientData ObjectType="Note"><x:MoveWithCells/><x:SizeWithCells/><x:Anchor>1, 15, 3, 2, 3, 31, 6, 1</x:Anchor><x:AutoFill>False</x:AutoFill><x:Row>3</x:Row><x:Column>1</x:Column></x:ClientData></v:shape></xml>'''
    workbook = entries["xl/workbook.xml"].decode()
    workbook = workbook.replace(
        "</workbook>",
        '<workbookProtection lockStructure="1"/>'
        '<definedNames><definedName name="CTOX_Amount">Overview!$B$2</definedName>'
        '<definedName name="CTOX_LocalInput" localSheetId="0">Overview!$B$3</definedName></definedNames>'
        '</workbook>',
    )
    entries["xl/workbook.xml"] = workbook.encode()
    content_types = entries["[Content_Types].xml"]
    content_types = content_types.replace(
        b"</Types>",
        b'<Default Extension="vml" ContentType="application/vnd.openxmlformats-officedocument.vmlDrawing"/><Override PartName="/xl/comments1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.comments+xml"/></Types>',
    )
    entries["[Content_Types].xml"] = content_types
    entries["docProps/core.xml"] = entries["docProps/core.xml"].replace(b"Spreadsheet Open Render Sheets", b"Spreadsheet Comments Names Protection")
    entries["customXml/ctox-spreadsheet-preserve.xml"] = b'''<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="spreadsheet.comments-names-protection">SPREADSHEET_COMMENTS_NAMES_PROTECTION_ESCROW_6E21</ctox-preserve>'''
    with ZipFile(SPREADSHEET_COMMENTS_NAMES_PROTECTION_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 12, 0, 0, 0)); info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_COMMENTS_NAMES_PROTECTION_OUTPUT)


def create_spreadsheet_charts() -> None:
    """Create a deterministic editable column-chart fixture with preserved parts."""
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    # Keep indices 6..17 valid for the inherited Details/Archive worksheets.
    values = [
        "Month", "Revenue", "January", "February", "March", "April",
        "CTOX DETAILS", "ORACLE_SHEET_DETAILS_24AF", "Item", "Quantity",
        "Consulting", "Support", "CTOX ARCHIVE", "ORACLE_SHEET_ARCHIVE_B83D",
        "Year", "Status", "2025", "Closed", "CTOX CHART SOURCE",
    ]
    entries["xl/sharedStrings.xml"] = (
        f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="{len(values)}" uniqueCount="{len(values)}">'''
        + "".join(f"<si><t>{value}</t></si>" for value in values)
        + "</sst>"
    ).encode()
    entries["xl/worksheets/sheet1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><dimension ref="A1:B6"/><sheetViews><sheetView workbookViewId="0"><selection activeCell="D2" sqref="D2"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="1" width="18" customWidth="1"/><col min="2" max="2" width="14" customWidth="1"/></cols><sheetData><row r="1"><c r="A1" t="s" s="2"><v>0</v></c><c r="B1" t="s" s="2"><v>1</v></c></row><row r="2"><c r="A2" t="s"><v>2</v></c><c r="B2"><v>120</v></c></row><row r="3"><c r="A3" t="s"><v>3</v></c><c r="B3"><v>185</v></c></row><row r="4"><c r="A4" t="s"><v>4</v></c><c r="B4"><v>160</v></c></row><row r="5"><c r="A5" t="s"><v>5</v></c><c r="B5"><v>240</v></c></row><row r="6"><c r="A6" t="s"><v>18</v></c></row></sheetData><drawing r:id="rId1"/></worksheet>'''
    entries["xl/worksheets/_rels/sheet1.xml.rels"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing" Target="../drawings/drawing1.xml"/></Relationships>'''
    entries["xl/drawings/drawing1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><xdr:wsDr xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><xdr:twoCellAnchor editAs="oneCell"><xdr:from><xdr:col>3</xdr:col><xdr:colOff>0</xdr:colOff><xdr:row>1</xdr:row><xdr:rowOff>0</xdr:rowOff></xdr:from><xdr:to><xdr:col>10</xdr:col><xdr:colOff>0</xdr:colOff><xdr:row>18</xdr:row><xdr:rowOff>0</xdr:rowOff></xdr:to><xdr:graphicFrame><xdr:nvGraphicFramePr><xdr:cNvPr id="2" name="CTOX Revenue Chart"/><xdr:cNvGraphicFramePr/></xdr:nvGraphicFramePr><xdr:xfrm><a:off x="2638440" y="250200"/><a:ext cx="3960000" cy="3238200"/></xdr:xfrm><a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"><c:chart xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" r:id="rId1"/></a:graphicData></a:graphic></xdr:graphicFrame><xdr:clientData/></xdr:twoCellAnchor></xdr:wsDr>'''
    entries["xl/drawings/_rels/drawing1.xml.rels"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart" Target="../charts/chart1.xml"/></Relationships>'''
    entries["xl/charts/chart1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><c:date1904 val="0"/><c:lang val="de-DE"/><c:roundedCorners val="0"/><c:chart><c:title><c:tx><c:rich><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="de-DE" sz="1400" b="1"/><a:t>CTOX Revenue 2026</a:t></a:r></a:p></c:rich></c:tx><c:layout/><c:overlay val="0"/></c:title><c:autoTitleDeleted val="0"/><c:plotArea><c:layout/><c:barChart><c:barDir val="col"/><c:grouping val="clustered"/><c:varyColors val="0"/><c:ser><c:idx val="0"/><c:order val="0"/><c:tx><c:strRef><c:f>Overview!$B$1</c:f><c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Revenue</c:v></c:pt></c:strCache></c:strRef></c:tx><c:spPr><a:solidFill><a:srgbClr val="176B5B"/></a:solidFill></c:spPr><c:cat><c:strRef><c:f>Overview!$A$2:$A$5</c:f><c:strCache><c:ptCount val="4"/><c:pt idx="0"><c:v>January</c:v></c:pt><c:pt idx="1"><c:v>February</c:v></c:pt><c:pt idx="2"><c:v>March</c:v></c:pt><c:pt idx="3"><c:v>April</c:v></c:pt></c:strCache></c:strRef></c:cat><c:val><c:numRef><c:f>Overview!$B$2:$B$5</c:f><c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="4"/><c:pt idx="0"><c:v>120</c:v></c:pt><c:pt idx="1"><c:v>185</c:v></c:pt><c:pt idx="2"><c:v>160</c:v></c:pt><c:pt idx="3"><c:v>240</c:v></c:pt></c:numCache></c:numRef></c:val></c:ser><c:dLbls><c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/><c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/></c:dLbls><c:gapWidth val="150"/><c:axId val="52743552"/><c:axId val="52747392"/></c:barChart><c:catAx><c:axId val="52743552"/><c:scaling><c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="b"/><c:tickLblPos val="nextTo"/><c:crossAx val="52747392"/><c:crosses val="autoZero"/><c:auto val="1"/><c:lblAlgn val="ctr"/><c:lblOffset val="100"/></c:catAx><c:valAx><c:axId val="52747392"/><c:scaling><c:orientation val="minMax"/></c:scaling><c:delete val="0"/><c:axPos val="l"/><c:numFmt formatCode="General" sourceLinked="1"/><c:majorGridlines/><c:tickLblPos val="nextTo"/><c:crossAx val="52743552"/><c:crosses val="autoZero"/><c:crossBetween val="between"/></c:valAx></c:plotArea><c:plotVisOnly val="1"/><c:dispBlanksAs val="gap"/><c:showDLblsOverMax val="0"/></c:chart><c:printSettings><c:headerFooter/><c:pageMargins b="0.75" l="0.7" r="0.7" t="0.75" header="0.3" footer="0.3"/><c:pageSetup/></c:printSettings></c:chartSpace>'''
    entries["[Content_Types].xml"] = entries["[Content_Types].xml"].replace(
        b"</Types>",
        b'<Override PartName="/xl/drawings/drawing1.xml" ContentType="application/vnd.openxmlformats-officedocument.drawing+xml"/><Override PartName="/xl/charts/chart1.xml" ContentType="application/vnd.openxmlformats-officedocument.drawingml.chart+xml"/></Types>',
    )
    entries["docProps/core.xml"] = entries["docProps/core.xml"].replace(
        b"Spreadsheet Open Render Sheets", b"Spreadsheet Charts"
    )
    entries["customXml/ctox-spreadsheet-preserve.xml"] = b'''<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="spreadsheet.charts">SPREADSHEET_CHART_ESCROW_53C8</ctox-preserve>'''
    with ZipFile(SPREADSHEET_CHARTS_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 12, 0, 0, 0)); info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_CHARTS_OUTPUT)


def create_spreadsheet_pivot_print_layout() -> None:
    """Create deterministic pivot-cache and print-layout OOXML for parity work."""
    with ZipFile(SPREADSHEET_OPEN_RENDER_OUTPUT, "r") as source:
        entries = {item.filename: source.read(item.filename) for item in source.infolist()}
    # Preserve inherited Details/Archive shared-string indices 6..17.
    values = [
        "Region", "Revenue", "North", "South", "CTOX PIVOT SOURCE", "PIVOT_PRINT_ORACLE_7C42",
        "CTOX DETAILS", "ORACLE_SHEET_DETAILS_24AF", "Item", "Quantity", "Consulting", "Support",
        "CTOX ARCHIVE", "ORACLE_SHEET_ARCHIVE_B83D", "Year", "Status", "2025", "Closed",
        "Product", "CTOX PIVOT REPORT", "Row Labels", "Sum of Revenue", "Grand Total",
    ]
    entries["xl/sharedStrings.xml"] = (
        f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="{len(values)}" uniqueCount="{len(values)}">'''
        + "".join(f"<si><t>{value}</t></si>" for value in values) + "</sst>"
    ).encode()
    entries["xl/worksheets/sheet1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheetPr><pageSetUpPr fitToPage="1"/></sheetPr><dimension ref="A1:F10"/><sheetViews><sheetView workbookViewId="0" view="pageBreakPreview"><selection activeCell="E3" sqref="E3"/></sheetView></sheetViews><sheetFormatPr defaultRowHeight="15"/><cols><col min="1" max="3" width="18" customWidth="1"/><col min="5" max="6" width="20" customWidth="1"/></cols><sheetData><row r="1"><c r="A1" t="s" s="2"><v>0</v></c><c r="B1" t="s" s="2"><v>18</v></c><c r="C1" t="s" s="2"><v>1</v></c></row><row r="2"><c r="A2" t="s"><v>2</v></c><c r="B2" t="s"><v>10</v></c><c r="C2"><v>420</v></c></row><row r="3"><c r="A3" t="s"><v>2</v></c><c r="B3" t="s"><v>11</v></c><c r="C3"><v>310</v></c><c r="E3" t="s" s="2"><v>19</v></c></row><row r="4"><c r="A4" t="s"><v>3</v></c><c r="B4" t="s"><v>10</v></c><c r="C4"><v>280</v></c><c r="E4" t="s"><v>20</v></c><c r="F4" t="s"><v>21</v></c></row><row r="5"><c r="E5" t="s"><v>2</v></c><c r="F5"><v>730</v></c></row><row r="6"><c r="E6" t="s"><v>3</v></c><c r="F6"><v>280</v></c></row><row r="7"><c r="E7" t="s"><v>22</v></c><c r="F7"><v>1010</v></c></row></sheetData><printOptions headings="1" gridLines="1" gridLinesSet="1" horizontalCentered="1"/><pageMargins left="0.5" right="0.5" top="0.75" bottom="0.75" header="0.3" footer="0.3"/><pageSetup paperSize="9" orientation="landscape" fitToWidth="1" fitToHeight="0" firstPageNumber="3" useFirstPageNumber="1" horizontalDpi="300" verticalDpi="300"/><headerFooter alignWithMargins="1"><oddHeader>&amp;LCTOX PIVOT&amp;CPRINT ORACLE&amp;RCONFIDENTIAL</oddHeader><oddFooter>&amp;LCTOX&amp;CPage &amp;P of &amp;N&amp;R2026</oddFooter></headerFooter><rowBreaks count="1" manualBreakCount="1"><brk id="7" min="0" max="16383" man="1"/></rowBreaks><colBreaks count="1" manualBreakCount="1"><brk id="3" min="0" max="1048575" man="1"/></colBreaks><pivotTableParts count="1"><pivotTablePart r:id="rIdPivot1"/></pivotTableParts></worksheet>'''
    entries["xl/worksheets/_rels/sheet1.xml.rels"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rIdPivot1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/pivotTable" Target="../pivotTables/pivotTable1.xml"/></Relationships>'''
    entries["xl/pivotTables/pivotTable1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><pivotTableDefinition xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" name="CTOXRevenuePivot" cacheId="1" dataCaption="Values" updatedVersion="8" minRefreshableVersion="3" useAutoFormatting="1" itemPrintTitles="1" createdVersion="8" indent="0"><location ref="E3:F7" firstHeaderRow="1" firstDataRow="2" firstDataCol="1"/><pivotFields count="3"><pivotField axis="axisRow" showAll="0"><items count="3"><item x="0"/><item x="1"/><item t="default"/></items></pivotField><pivotField showAll="0"/><pivotField dataField="1" showAll="0"/></pivotFields><rowFields count="1"><field x="0"/></rowFields><rowItems count="3"><i><x v="0"/></i><i><x v="1"/></i><i t="grand"><x/></i></rowItems><colFields count="1"><field x="-2"/></colFields><colItems count="1"><i><x v="0"/></i></colItems><dataFields count="1"><dataField name="Sum of Revenue" fld="2" subtotal="sum"/></dataFields><pivotTableStyleInfo name="PivotStyleMedium9" showRowHeaders="1" showColHeaders="1" showRowStripes="0" showColStripes="0" showLastColumn="0"/></pivotTableDefinition>'''
    entries["xl/pivotTables/_rels/pivotTable1.xml.rels"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/pivotCacheDefinition" Target="../pivotCache/pivotCacheDefinition1.xml"/></Relationships>'''
    entries["xl/pivotCache/pivotCacheDefinition1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><pivotCacheDefinition xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" r:id="rId1" saveData="1" refreshOnLoad="0" recordCount="3" createdVersion="8" refreshedVersion="8" minRefreshableVersion="3"><cacheSource type="worksheet"><worksheetSource ref="A1:C4" sheet="Overview"/></cacheSource><cacheFields count="3"><cacheField name="Region" numFmtId="0"><sharedItems count="2"><s v="North"/><s v="South"/></sharedItems></cacheField><cacheField name="Product" numFmtId="0"><sharedItems count="2"><s v="Consulting"/><s v="Support"/></sharedItems></cacheField><cacheField name="Revenue" numFmtId="0"><sharedItems containsString="0" containsNumber="1" containsInteger="1" minValue="280" maxValue="420"/></cacheField></cacheFields></pivotCacheDefinition>'''
    entries["xl/pivotCache/_rels/pivotCacheDefinition1.xml.rels"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/pivotCacheRecords" Target="pivotCacheRecords1.xml"/></Relationships>'''
    entries["xl/pivotCache/pivotCacheRecords1.xml"] = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?><pivotCacheRecords xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="3"><r><x v="0"/><x v="0"/><n v="420"/></r><r><x v="0"/><x v="1"/><n v="310"/></r><r><x v="1"/><x v="0"/><n v="280"/></r></pivotCacheRecords>'''
    workbook = entries["xl/workbook.xml"].decode().replace(
        "</workbook>",
        '<pivotCaches><pivotCache cacheId="1" r:id="rIdPivotCache1"/></pivotCaches><definedNames><definedName name="_xlnm.Print_Area" localSheetId="0">Overview!$A$1:$F$10</definedName><definedName name="_xlnm.Print_Titles" localSheetId="0">Overview!$1:$1</definedName></definedNames></workbook>',
    )
    entries["xl/workbook.xml"] = workbook.encode()
    entries["xl/_rels/workbook.xml.rels"] = entries["xl/_rels/workbook.xml.rels"].replace(
        b"</Relationships>",
        b'<Relationship Id="rIdPivotCache1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/pivotCacheDefinition" Target="pivotCache/pivotCacheDefinition1.xml"/></Relationships>',
    )
    entries["[Content_Types].xml"] = entries["[Content_Types].xml"].replace(
        b"</Types>",
        b'<Override PartName="/xl/pivotTables/pivotTable1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.pivotTable+xml"/><Override PartName="/xl/pivotCache/pivotCacheDefinition1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.pivotCacheDefinition+xml"/><Override PartName="/xl/pivotCache/pivotCacheRecords1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.pivotCacheRecords+xml"/></Types>',
    )
    entries["docProps/core.xml"] = entries["docProps/core.xml"].replace(b"Spreadsheet Open Render Sheets", b"Spreadsheet Pivot Print Layout")
    entries["customXml/ctox-spreadsheet-preserve.xml"] = b'''<?xml version="1.0" encoding="UTF-8"?><ctox-preserve feature="spreadsheet.pivot-print-layout">SPREADSHEET_PIVOT_PRINT_ESCROW_7C42</ctox-preserve>'''
    with ZipFile(SPREADSHEET_PIVOT_PRINT_LAYOUT_OUTPUT, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data in sorted(entries.items()):
            info = ZipInfo(name, date_time=(2026, 7, 13, 0, 0, 0)); info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    print(SPREADSHEET_PIVOT_PRINT_LAYOUT_OUTPUT)


def main() -> None:
    create_open_render_zoom()
    create_edit_save()
    create_undo_clipboard_keyboard()
    create_character_paragraph_formatting()
    create_styles_lists_numbering()
    create_tables()
    create_images_positioning()
    create_sections_headers_footers()
    create_links_bookmarks_fields()
    create_comments_track_changes()
    create_drawings_charts()
    create_spreadsheet_open_render_sheets()
    create_spreadsheet_edit_save()
    create_spreadsheet_undo_clipboard_fill()
    create_spreadsheet_cell_format_rows_columns()
    create_spreadsheet_formulas_references()
    create_spreadsheet_multi_sheet_merge_freeze()
    create_spreadsheet_sort_filter_tables()
    create_spreadsheet_validation_conditional_formatting()
    create_spreadsheet_comments_names_protection()
    create_spreadsheet_charts()
    create_spreadsheet_pivot_print_layout()


def canonicalize_zip(path: Path) -> None:
    temporary = path.with_suffix(".canonical.docx")
    with ZipFile(path, "r") as source:
        entries = [
            (item.filename, source.read(item.filename), item.external_attr)
            for item in source.infolist()
        ]
    with ZipFile(temporary, "w", compression=ZIP_DEFLATED, compresslevel=9) as target:
        for name, data, external_attr in sorted(entries):
            info = ZipInfo(name, date_time=(2026, 7, 10, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            info.external_attr = external_attr
            target.writestr(info, data)
    temporary.replace(path)


if __name__ == "__main__":
    main()
